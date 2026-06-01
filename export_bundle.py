"""Build ZIP export: HTML, audit JSON, summary CSV, PDF snapshot.

See docs/FOLDER_MAP.md."""

from __future__ import annotations

import csv
import io
import json
import os
import zipfile
from datetime import datetime
from typing import Any

import pandas as pd


AR_NUM_TRANSLATION = str.maketrans("٠١٢٣٤٥٦٧٨٩٫٬", "0123456789.,")


def _first_series(df: pd.DataFrame, col: str) -> pd.Series:
    obj = df[col]
    if isinstance(obj, pd.DataFrame):
        return obj.iloc[:, 0]
    return obj


def _to_numeric_series(series: pd.Series) -> pd.Series:
    """Parse numbers robustly, including Arabic-Indic digits and separators."""
    if pd.api.types.is_numeric_dtype(series):
        return pd.to_numeric(series, errors="coerce")
    s = series.astype(str).str.strip().str.translate(AR_NUM_TRANSLATION)
    s = s.str.replace("%", "", regex=False).str.replace(",", "", regex=False)
    s = s.str.replace(r"[^\d\.\-\+]", "", regex=True)
    return pd.to_numeric(s, errors="coerce")


def _pick_metric_column(df: pd.DataFrame, audit: dict[str, Any]) -> str | None:
    detected = (audit.get("detected_columns") or {}) if isinstance(audit, dict) else {}
    for k in ("revenue", "sales", "profit", "cost"):
        c = detected.get(k)
        if c in df.columns:
            return str(c)
    num_cols = list(df.select_dtypes(include=["number"]).columns)
    if num_cols:
        return str(num_cols[0])
    best_col: str | None = None
    best_score = -1
    for col in df.columns:
        ser = _to_numeric_series(_first_series(df, col))
        score = int(ser.notna().sum())
        if score > best_score and score >= max(2, len(df) // 3):
            best_col = str(col)
            best_score = score
    return best_col


def _pick_period_column(df: pd.DataFrame, audit: dict[str, Any]) -> str | None:
    detected = (audit.get("detected_columns") or {}) if isinstance(audit, dict) else {}
    c = detected.get("period")
    if c in df.columns:
        return str(c)
    name_hits = ("date", "month", "period", "quarter", "year", "time", "تاريخ", "شهر", "سنة")
    for col in df.columns:
        s = str(col).lower()
        if any(x in s for x in name_hits):
            return str(col)
    for col in df.columns:
        ser = pd.to_datetime(_first_series(df, col), errors="coerce", format="mixed")
        if int(ser.notna().sum()) >= max(2, len(df) // 2):
            return str(col)
    return None


def _pick_segment_column(df: pd.DataFrame, audit: dict[str, Any], metric_col: str | None) -> str | None:
    detected = (audit.get("detected_columns") or {}) if isinstance(audit, dict) else {}
    c = detected.get("segment")
    if c in df.columns:
        return str(c)
    best: str | None = None
    best_n = 0
    for col in df.columns:
        if str(col) == str(metric_col):
            continue
        n = int(df[col].nunique(dropna=True))
        if 2 <= n <= 30 and n > best_n:
            best = str(col)
            best_n = n
    return best


def _sorted_period_labels(values: list[str]) -> list[str]:
    if not values:
        return values
    ser = pd.Series(values, dtype="string")
    dt = pd.to_datetime(ser, errors="coerce", format="mixed")
    if int(dt.notna().sum()) >= max(2, len(values) // 2):
        tmp = pd.DataFrame({"label": values, "_dt": dt})
        tmp["_dt"] = tmp["_dt"].fillna(pd.Timestamp.max)
        tmp = tmp.sort_values("_dt")
        return [str(x) for x in tmp["label"].tolist()]
    return sorted(values)


def _row_order_series(metric: pd.Series, max_points: int = 24) -> tuple[list[str], list[float]]:
    m = metric.fillna(0.0).reset_index(drop=True)
    if len(m) == 0:
        return [], []
    chunk = max(1, int((len(m) + max_points - 1) // max_points))
    labels: list[str] = []
    vals: list[float] = []
    for i in range(0, len(m), chunk):
        j = min(i + chunk, len(m))
        labels.append(f"#{i + 1}-#{j}")
        vals.append(float(m.iloc[i:j].sum()))
    return labels, vals


def build_summary_pptx(df: pd.DataFrame, audit: dict[str, Any]) -> bytes:
    """PowerPoint analysis deck with overview + real trend/segment charts."""
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches, Pt
    except ImportError:
        return b""

    prs = Presentation()
    metric_col = _pick_metric_column(df, audit)
    period_col = _pick_period_column(df, audit)
    segment_col = _pick_segment_column(df, audit, metric_col)
    metric_series = (
        _to_numeric_series(_first_series(df, metric_col)).dropna()
        if metric_col and metric_col in df.columns
        else pd.Series(dtype=float)
    )

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Analysis Export Snapshot"
    subtitle = slide.placeholders[1]
    subtitle.text = (
        f"Source: {audit.get('source_name', '')}\n"
        f"Generated: {audit.get('generated_at', datetime.now().isoformat())}\n"
        f"Report ID: {audit.get('report_id', '')}"
    )

    # Slide 2: Executive overview (computed from actual data)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Overview"
    body = slide.placeholders[1].text_frame
    body.clear()
    lines: list[str] = [
        f"Rows x Columns: {audit.get('rows', df.shape[0])} x {audit.get('columns', df.shape[1])}",
        f"Missing cells: {audit.get('missing_cells', int(df.isna().sum().sum()))}",
        f"Duplicate rows: {audit.get('duplicate_rows', int(df.duplicated().sum()))}",
    ]
    if metric_col and metric_col in df.columns:
        m = metric_series
        if len(m) > 0:
            lines.extend(
                [
                    f"Primary metric: {metric_col}",
                    f"Total: {float(m.sum()):,.2f}",
                    f"Average: {float(m.mean()):,.2f}",
                    f"Range: {float(m.min()):,.2f} to {float(m.max()):,.2f}",
                ]
            )
    if period_col:
        lines.append(f"Period dimension: {period_col}")
    else:
        lines.append("Period dimension: row-order fallback")
    if segment_col:
        lines.append(f"Segment dimension: {segment_col}")
    else:
        lines.append("Segment dimension: overall fallback")
    for i, line in enumerate(lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = str(line)
        p.level = 0
        p.font.size = Pt(14)

    # Slide 3: Trend chart
    if metric_col and metric_col in df.columns and len(metric_series) > 0:
        labels: list[str] = []
        vals: list[float] = []
        trend_title: str
        if period_col and period_col in df.columns:
            w = pd.DataFrame(
                {
                    "_p": _first_series(df, period_col).astype(str).str.strip().replace({"": "Unknown"}),
                    "_m": _to_numeric_series(_first_series(df, metric_col)).fillna(0.0),
                }
            )
            grp = w.groupby("_p", dropna=False)["_m"].sum()
            labels = _sorted_period_labels([str(x) for x in grp.index.tolist()])
            vals = [float(grp.get(l, 0.0)) for l in labels]
            trend_title = f"Trend Analysis ({metric_col} by {period_col})"
        else:
            labels, vals = _row_order_series(metric_series, max_points=24)
            trend_title = f"Trend Analysis ({metric_col} by row order)"
        if labels and vals:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
            slide.shapes.title.text = trend_title
            chart_data = CategoryChartData()
            chart_data.categories = labels[:24]
            chart_data.add_series(metric_col, vals[:24])
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE_MARKERS,
                Inches(0.7),
                Inches(1.5),
                Inches(12.0),
                Inches(5.0),
                chart_data,
            ).chart
            chart.has_legend = False

    # Slide 4: Segment chart
    if metric_col and metric_col in df.columns and len(metric_series) > 0:
        labels: list[str] = []
        vals: list[float] = []
        seg_title: str
        if segment_col and segment_col in df.columns:
            w = pd.DataFrame(
                {
                    "_s": _first_series(df, segment_col).astype(str).str.strip().replace({"": "Unknown"}),
                    "_m": _to_numeric_series(_first_series(df, metric_col)).fillna(0.0),
                }
            )
            grp = (
                w.groupby("_s", dropna=False)["_m"]
                .sum()
                .sort_values(ascending=False)
                .head(10)
            )
            labels = [str(x) for x in grp.index.tolist()]
            vals = [float(x) for x in grp.values.tolist()]
            seg_title = f"Segment Analysis (Top {len(labels)} by {metric_col})"
        else:
            labels = ["All data"]
            vals = [float(metric_series.sum())]
            seg_title = f"Segment Analysis (overall {metric_col})"
        if labels and vals:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
            slide.shapes.title.text = seg_title
            chart_data = CategoryChartData()
            chart_data.categories = labels
            chart_data.add_series(metric_col, vals)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(0.7),
                Inches(1.5),
                Inches(12.0),
                Inches(5.0),
                chart_data,
            ).chart
            chart.has_legend = False

    # Slide 5: KPI + anomalies
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Finance KPIs and Alerts"
    body = slide.placeholders[1].text_frame
    body.clear()
    kpis = audit.get("finance_kpis") or []
    anomalies = audit.get("triggered_anomalies") or []
    lines = []
    if kpis:
        lines.append("KPIs:")
        for row in kpis[:10]:
            lines.append(f"- {row.get('name', '')}: {row.get('value', '')}")
    else:
        lines.append("KPIs: (none)")
    if anomalies:
        lines.append("Alerts:")
        for a in anomalies[:12]:
            lines.append(f"- [{a.get('severity', '')}] {a.get('rule', '')} (n={a.get('count', '')})")
    else:
        lines.append("Alerts: (none triggered)")
    for i, line in enumerate(lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(14)

    # Slide 6: Data sample table
    sample = df.head(12)
    cols = list(sample.columns[:6])
    rows = len(sample) + 1
    if rows >= 2 and cols:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
        slide.shapes.title.text = "Data Sample (first rows)"
        table_shape = slide.shapes.add_table(
            rows, len(cols), Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4)
        )
        table = table_shape.table
        for c, col in enumerate(cols):
            table.cell(0, c).text = str(col)
        for r_idx, (_, row) in enumerate(sample.iterrows(), start=1):
            for c, col in enumerate(cols):
                v = row[col]
                table.cell(r_idx, c).text = "" if pd.isna(v) else str(v)

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()


def build_summary_csv(df: pd.DataFrame, audit: dict[str, Any]) -> bytes:
    """UTF-8 BOM CSV: metadata, detected columns, KPIs, anomalies, numeric profile, data sample."""
    buf = io.StringIO()
    w = csv.writer(buf)

    w.writerow(["# Metadata"])
    w.writerow(["key", "value"])
    meta_keys = (
        "report_id",
        "report_version",
        "generated_at",
        "source_name",
        "sheet_name",
        "mode",
        "rows",
        "columns",
        "missing_cells",
        "duplicate_rows",
        "content_sha256",
    )
    for k in meta_keys:
        if k in audit and audit[k] is not None:
            w.writerow([k, audit[k]])

    w.writerow([])
    w.writerow(["# Detected columns"])
    w.writerow(["role", "column_name"])
    dc = audit.get("detected_columns") or {}
    if isinstance(dc, dict):
        for role, col in dc.items():
            w.writerow([role, col if col is not None else ""])

    w.writerow([])
    w.writerow(["# Schema warnings"])
    w.writerow(["warning"])
    for x in audit.get("schema_warnings") or []:
        w.writerow([x])

    w.writerow([])
    w.writerow(["# Finance KPIs"])
    w.writerow(["name", "value"])
    for row in audit.get("finance_kpis") or []:
        w.writerow([row.get("name", ""), row.get("value", "")])

    w.writerow([])
    w.writerow(["# Anomaly rules"])
    w.writerow(["severity", "rule", "count"])
    for a in audit.get("triggered_anomalies") or []:
        w.writerow([a.get("severity", ""), a.get("rule", ""), a.get("count", "")])

    w.writerow([])
    w.writerow(["# Numeric column statistics"])
    num_df = df.select_dtypes(include=["number"])
    if num_df.shape[1] > 0:
        buf.write(num_df.describe().to_csv(lineterminator="\n"))
    else:
        w.writerow(["(no numeric columns)"])

    w.writerow([])
    w.writerow(["# Data sample (first 200 rows)"])
    sample = df.head(200)
    buf.write(sample.to_csv(index=False, lineterminator="\n"))

    return buf.getvalue().encode("utf-8-sig")


def _fpdf_dejavu_path() -> str | None:
    try:
        import fpdf

        font_dir = os.path.join(os.path.dirname(fpdf.__file__), "font")
        for name in ("DejaVuSans.ttf", "dejavusans.ttf"):
            p = os.path.join(font_dir, name)
            if os.path.isfile(p):
                return p
    except Exception:
        pass
    return None


def build_summary_pdf(df: pd.DataFrame, audit: dict[str, Any]) -> bytes:
    """One-page PDF snapshot for archiving (Unicode via DejaVu if available)."""
    try:
        from fpdf import FPDF
    except ImportError:
        return b""

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    font_path = _fpdf_dejavu_path()
    if font_path:
        pdf.add_font("DejaVu", "", font_path)
        pdf.set_font("DejaVu", size=10)
        font = "DejaVu"
    else:
        pdf.set_font("Helvetica", size=10)
        font = "Helvetica"

    pdf.add_page()
    epw = pdf.epw
    pdf.set_font(font, "B", 16)
    pdf.multi_cell(epw, 8, "Analysis export snapshot", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font(font, size=9)
    pdf.multi_cell(
        epw,
        5,
        f"Generated: {audit.get('generated_at', datetime.now().isoformat())}",
        new_x="LMARGIN",
        new_y="NEXT",
    )
    pdf.ln(3)

    lines = [
        f"Report ID: {audit.get('report_id', '')}",
        f"Source: {audit.get('source_name', '')}",
        f"Version: {audit.get('report_version', '')}",
        f"SHA-256: {audit.get('content_sha256', '')}",
        f"Shape: {audit.get('rows', df.shape[0])} rows x {audit.get('columns', df.shape[1])} columns",
        f"Missing cells: {audit.get('missing_cells', '')} | Duplicate rows: {audit.get('duplicate_rows', '')}",
    ]
    for line in lines:
        pdf.multi_cell(epw, 5, _safe_pdf_text(line, font == "Helvetica"))
    pdf.ln(2)

    pdf.set_font(font, "B", 11)
    pdf.multi_cell(epw, 7, "Finance KPIs", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font(font, size=9)
    for row in audit.get("finance_kpis") or []:
        pdf.multi_cell(
            epw,
            5,
            _safe_pdf_text(f"  • {row.get('name', '')}: {row.get('value', '')}", font == "Helvetica"),
        )

    pdf.ln(2)
    pdf.set_font(font, "B", 11)
    pdf.multi_cell(epw, 7, "Anomalies", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font(font, size=9)
    alerts = audit.get("triggered_anomalies") or []
    if not alerts:
        pdf.multi_cell(epw, 5, "  (none triggered)")
    else:
        for a in alerts:
            pdf.multi_cell(
                epw,
                5,
                _safe_pdf_text(
                    f"  [{a.get('severity', '')}] {a.get('rule', '')} (n={a.get('count', '')})",
                    font == "Helvetica",
                ),
            )

    pdf.ln(2)
    pdf.set_font(font, "I", 8)
    pdf.multi_cell(
        epw,
        4,
        _safe_pdf_text(
            "Full detail: open report.html in this bundle. Correlations and charts are in the HTML only.",
            font == "Helvetica",
        ),
    )

    out = pdf.output(dest="S")
    if isinstance(out, str):
        return out.encode("latin-1", errors="replace")
    return bytes(out)


def _safe_pdf_text(s: str, ascii_only: bool) -> str:
    if not ascii_only:
        return s
    return s.encode("ascii", errors="replace").decode("ascii")


def build_readme() -> str:
    return """Analysis export bundle
========================

Contents:
  report.html   — Full interactive dashboard (open in a browser).
  audit.json    — Audit trail and detected columns (machine-readable).
  summary.csv   — Metadata, KPIs, anomalies, numeric stats, and data sample.
  summary.pdf   — One-page PDF snapshot (if generated; requires fpdf2).
  summary.pptx  — PowerPoint snapshot (if generated; requires python-pptx).

Keep this ZIP with your records. The SHA-256 in audit.json ties the report to the input data.
"""


def create_export_zip(
    html: str,
    audit: dict[str, Any],
    df: pd.DataFrame,
    *,
    readme_extra: str = "",
) -> bytes:
    """Return ZIP file bytes."""
    bio = io.BytesIO()
    csv_bytes = build_summary_csv(df, audit)
    pdf_bytes = build_summary_pdf(df, audit)
    pptx_bytes = build_summary_pptx(df, audit)

    with zipfile.ZipFile(bio, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("report.html", html)
        zf.writestr(
            "audit.json",
            json.dumps(audit, ensure_ascii=False, indent=2),
        )
        zf.writestr("summary.csv", csv_bytes)
        if pdf_bytes:
            zf.writestr("summary.pdf", pdf_bytes)
        if pptx_bytes:
            zf.writestr("summary.pptx", pptx_bytes)
        readme = build_readme() + ("\n" + readme_extra if readme_extra else "")
        zf.writestr("README.txt", readme)

    bio.seek(0)
    return bio.getvalue()
