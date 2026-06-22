"""PPTX export for legal text details."""
from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt


def build_legal_text_pptx(
    legal_text: str,
    fields: list[dict[str, str]] | None = None,
    *,
    title: str = "النص النظامي",
) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(9)

    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True

    body_top = Inches(1.1)
    body_box = slide.shapes.add_textbox(left, body_top, width, Inches(1.2))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_tf.text = legal_text or "—"
    body_tf.paragraphs[0].font.size = Pt(14)

    y = Inches(2.5)
    for field in fields or []:
        label = str(field.get("label", "")).strip()
        value = str(field.get("value", "")).strip()
        if not label and not value:
            continue
        box = slide.shapes.add_textbox(left, y, width, Inches(0.45))
        p = box.text_frame.paragraphs[0]
        p.text = f"{label}: {value}"
        p.font.size = Pt(12)
        y += Inches(0.5)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
